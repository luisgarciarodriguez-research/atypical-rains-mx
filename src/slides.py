"""
Generador de presentación PPTX — Hallazgos iniciales del proyecto.

Produce una presentación en formato PowerPoint (16:9, 33.87 × 19.05 cm) con
python-pptx, dirigida a la exposición de avances ante el comité tutorial.
Incluye las siguientes diapositivas:

  1  Portada con logotipo UNAM/IIMAS y datos del proyecto.
  2  Motivación y alcance: precipitación atípica en México 2013–2026.
  3  Fuente de datos: red SMN, 1 959 estaciones, 161 meses.
  4  Calidad de datos: mapa de completitud y diagrama MCAR.
  5  Detección de anomalías: esquema de 4 capas y resultados de consenso.
  6  Análisis composicional (CoDA): símplex, CLR/ILR y SBP climatológica.
  7  Regímenes pluviométricos: mapa de clusters y perfiles medios por régimen.
  8  Conclusiones y trabajo futuro.

Todas las figuras se importan desde ``outputs/figures/`` en formato PNG.
La salida se escribe en ``outputs/reports/hallazgos_iniciales_lluvias_mx.pptx``.
Punto de entrada: ``build_slides(verbose=True)``.

Proyecto : Detección de Precipitación Atípica en México (2013–2026)
Unidad   : Universidad Nacional Autónoma de México (UNAM)
           Instituto de Investigaciones en Matemáticas Aplicadas y en Sistemas (IIMAS)
           Posgrado en Ciencia e Ingeniería de la Computación (PCIC)
Investigador líder : Dr. José Antonio Neme Castillo
Contribuidor       : Luis García Rodríguez
"""

from __future__ import annotations

from datetime import date
from io import BytesIO
from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Pt

from src.config import FIGURES, ROOT

OUTPUT_PPTX = ROOT / "outputs" / "reports" / "hallazgos_iniciales_lluvias_mx.pptx"

# ── Dimensiones (16:9, 33.87 × 19.05 cm) ─────────────────────────────────────
SW = Cm(33.867)
SH = Cm(19.05)

# Coordenadas del layout base
BAR_H   = Cm(2.4)          # alto de la barra de título
MARG_L  = Cm(1.6)          # margen izquierdo de contenido
MARG_R  = Cm(1.6)
CONT_Y  = BAR_H + Cm(0.3)  # inicio del área de contenido
CONT_H  = SH - CONT_Y - Cm(0.5)

# ── Paleta ────────────────────────────────────────────────────────────────────
AZUL_OSC  = RGBColor(0x15, 0x65, 0xC0)
AZUL_MED  = RGBColor(0x1E, 0x88, 0xE5)
AZUL_CLAR = RGBColor(0xBB, 0xDE, 0xFB)
GRIS_OSC  = RGBColor(0x42, 0x42, 0x42)
GRIS_MED  = RGBColor(0x75, 0x75, 0x75)
GRIS_CLAR = RGBColor(0xF5, 0xF5, 0xF5)
BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO     = RGBColor(0x00, 0x00, 0x00)
VERDE     = RGBColor(0x2E, 0x7D, 0x32)
ROJO      = RGBColor(0xC6, 0x28, 0x28)
NARANJA   = RGBColor(0xE6, 0x51, 0x00)

# ── Helpers de bajo nivel ─────────────────────────────────────────────────────

def _add_rect(slide, x, y, w, h, fill: RGBColor | None = None,
              line_color: RGBColor | None = None, line_w: float = 0):
    """Agrega un rectángulo (sin texto) a la diapositiva."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text: str, x, y, w, h,
              size: int = 18, bold: bool = False, italic: bool = False,
              color: RGBColor = NEGRO, align=PP_ALIGN.LEFT,
              wrap: bool = True, font: str = "Calibri") -> None:
    """Agrega un cuadro de texto."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_para(tf, text: str, size: int = 14, bold: bool = False,
              color: RGBColor = GRIS_OSC, align=PP_ALIGN.LEFT,
              space_before: float = 0, bullet: str = ""):
    """Agrega un párrafo con run a un TextFrame existente."""
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = _Pt(space_before)
    run = p.add_run()
    run.text = (bullet + " " if bullet else "") + text
    run.font.name  = "Calibri"
    run.font.size  = _Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return p


def _add_img(slide, path: Path, x, y, w, max_h=None):
    """Inserta imagen respetando aspect ratio."""
    if not path.exists():
        _add_text(slide, f"[{path.name}]", x, y, w, Cm(2),
                  size=9, color=GRIS_MED)
        return
    with PILImage.open(path) as pil:
        iw, ih = pil.size
    h = int(w * ih / iw)
    if max_h and h > max_h:
        h = max_h
        w = int(h * iw / ih)
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def _title_bar(slide, title: str, subtitle: str = ""):
    """Barra azul superior con título y subtítulo opcional."""
    _add_rect(slide, 0, 0, SW, BAR_H, fill=AZUL_OSC)
    _add_text(slide, title,
              MARG_L, Cm(0.25), SW - MARG_L * 2, Cm(1.4),
              size=22, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        _add_text(slide, subtitle,
                  MARG_L, Cm(1.5), SW - MARG_L * 2, Cm(0.8),
                  size=13, color=AZUL_CLAR, align=PP_ALIGN.LEFT)


def _section_label(slide, label: str, x, y, w=Cm(5)):
    """Etiqueta de sección pequeña."""
    _add_rect(slide, x, y, w, Cm(0.55), fill=AZUL_CLAR)
    _add_text(slide, label, x + Cm(0.2), y + Cm(0.05),
              w - Cm(0.4), Cm(0.5),
              size=9, bold=True, color=AZUL_OSC, align=PP_ALIGN.LEFT)


def _divider(slide, y, color: RGBColor = AZUL_CLAR):
    """Línea divisoria horizontal."""
    line = slide.shapes.add_shape(1, MARG_L, y, SW - MARG_L * 2, Cm(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _stat_box(slide, x, y, w, h, value: str, label: str,
              bg: RGBColor = AZUL_CLAR, val_color: RGBColor = AZUL_OSC):
    """Caja de estadístico grande (value arriba, label abajo)."""
    _add_rect(slide, x, y, w, h, fill=bg)
    _add_text(slide, value, x, y + Cm(0.25), w, h * 0.55,
              size=26, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    _add_text(slide, label, x, y + h * 0.55, w, h * 0.4,
              size=10, bold=False, color=GRIS_OSC, align=PP_ALIGN.CENTER, wrap=True)


# ── Diapositivas ──────────────────────────────────────────────────────────────

def slide_portada(prs: Presentation, today: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Fondo degradado simulado: bloque azul oscuro superior
    _add_rect(slide, 0, 0, SW, SH * 0.55, fill=AZUL_OSC)
    _add_rect(slide, 0, SH * 0.55, SW, SH * 0.45, fill=BLANCO)

    # Franja decorativa
    _add_rect(slide, 0, SH * 0.55, SW, Cm(0.18), fill=AZUL_MED)

    # Logotipo textual UNAM
    _add_text(slide, "UNAM  ·  DCIC  ·  Estadística Climática",
              MARG_L, Cm(0.6), SW - MARG_L * 2, Cm(0.8),
              size=11, color=AZUL_CLAR, align=PP_ALIGN.CENTER)

    # Título principal
    _add_text(slide, "Análisis de Precipitación\nAtípica en México",
              Cm(2), Cm(1.8), SW - Cm(4), Cm(5),
              size=38, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, wrap=True)

    # Subtítulo
    _add_text(slide, "Detección de Anomalías y Clasificación de Regímenes Pluviométricos\n2013 – 2026",
              Cm(3), Cm(6.5), SW - Cm(6), Cm(2),
              size=16, color=AZUL_CLAR, align=PP_ALIGN.CENTER, wrap=True)

    # Metadatos
    _add_text(slide, f"Hallazgos Iniciales  ·  {today}",
              Cm(3), SH * 0.55 + Cm(0.8), SW - Cm(6), Cm(0.8),
              size=12, color=GRIS_MED, align=PP_ALIGN.CENTER)

    # Cuatro cifras clave
    stats = [
        ("1,959", "estaciones\npluviométricas"),
        ("161", "meses\n(2013–2026)"),
        ("14,290", "anomalías\nconfirmadas"),
        ("14", "regímenes\nidentificados"),
    ]
    box_w = (SW - MARG_L * 2 - Cm(0.8)) / 4
    box_h = Cm(2.4)
    bx    = MARG_L
    by    = SH * 0.55 + Cm(2.0)
    for val, lbl in stats:
        _stat_box(slide, bx, by, box_w - Cm(0.2), box_h, val, lbl,
                  bg=GRIS_CLAR, val_color=AZUL_OSC)
        bx += box_w


def slide_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Agenda", "Estructura de la presentación")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    items = [
        ("01", "Dataset y Exploración",
         "1,959 estaciones · 161 meses · 46.2% faltante"),
        ("02", "Detección de Anomalías",
         "4 capas independientes · 14,290 anomalías confirmadas"),
        ("03", "Subconjunto Analítico (CoDA)",
         "1,302 estaciones · reemplazo multiplicativo de ceros"),
        ("04", "Transformaciones Log-Ratio",
         "CLR · ILR-SBP climatológico · verificación isométrica"),
        ("05", "Regímenes Pluviométricos",
         "K=14 clusters · 3 métodos · mapas y perfiles composicionales"),
    ]
    start_y = CONT_Y + Cm(0.3)
    item_h  = Cm(2.7)
    for i, (num, title, desc) in enumerate(items):
        y = start_y + i * item_h
        # Círculo numerado
        _add_rect(slide, MARG_L, y + Cm(0.3), Cm(1.1), Cm(1.1), fill=AZUL_OSC)
        _add_text(slide, num, MARG_L, y + Cm(0.3), Cm(1.1), Cm(1.1),
                  size=15, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        # Título del item
        _add_text(slide, title,
                  MARG_L + Cm(1.4), y + Cm(0.2), SW - MARG_L * 2 - Cm(1.4), Cm(0.9),
                  size=15, bold=True, color=AZUL_OSC)
        _add_text(slide, desc,
                  MARG_L + Cm(1.4), y + Cm(1.0), SW - MARG_L * 2 - Cm(1.4), Cm(0.7),
                  size=11, color=GRIS_MED)
        if i < len(items) - 1:
            _divider(slide, y + item_h - Cm(0.15), GRIS_CLAR)


def slide_dataset(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Dataset y Cobertura", "SMN · 1,959 estaciones pluviométricas")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    # Mapa izquierda
    map_w = Cm(18)
    _add_img(slide, FIGURES / "T1.1.3_station_coverage.png",
             MARG_L, CONT_Y + Cm(0.2), map_w, max_h=Cm(13.5))

    # Stats derecha
    rx = MARG_L + map_w + Cm(0.5)
    rw = SW - rx - MARG_R
    bh = Cm(2.5)
    gap = Cm(0.4)

    stats_r = [
        ("1,959", "Estaciones", AZUL_OSC, AZUL_CLAR),
        ("161",   "Meses analizados", AZUL_OSC, AZUL_CLAR),
        ("46.2%", "Datos faltantes", ROJO, RGBColor(0xFF, 0xCC, 0xBC)),
        ("32",    "Estados del país", VERDE, RGBColor(0xC8, 0xE6, 0xC9)),
    ]
    by = CONT_Y + Cm(0.5)
    for val, lbl, vc, bg in stats_r:
        _stat_box(slide, rx, by, rw, bh, val, lbl, bg=bg, val_color=vc)
        by += bh + gap

    by += Cm(0.2)
    _add_text(slide, "Código de faltante: −99.0\nPatrón no aleatorio (dropout operativo\nen noreste, cierre de estaciones)",
              rx, by, rw, Cm(2.5), size=10, color=GRIS_MED, wrap=True)


def slide_faltantes(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Análisis de Datos Faltantes",
               "Patrones espaciales y temporales de omisión")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    fig_w = (SW - MARG_L * 2 - Cm(0.5)) / 2
    _add_img(slide, FIGURES / "missing_by_state_year.png",
             MARG_L, CONT_Y + Cm(0.2), fig_w, max_h=Cm(12.5))

    rx = MARG_L + fig_w + Cm(0.5)
    _add_img(slide, FIGURES / "missing_dropout.png",
             rx, CONT_Y + Cm(0.2), fig_w, max_h=Cm(12.5))

    # Caption
    _add_text(slide,
              "Izq.: % faltante por estado y año  ·  Der.: curva de dropout temporal por estación",
              MARG_L, SH - Cm(1.3), SW - MARG_L * 2, Cm(0.8),
              size=9, color=GRIS_MED, align=PP_ALIGN.CENTER)


def slide_distribucion(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Distribución de Precipitación",
               "Marcada asimetría positiva · Patrón monzónico Jun–Sep")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    fig_w = (SW - MARG_L * 2 - Cm(0.5)) / 2
    _add_img(slide, FIGURES / "dist_seasonal_boxplot.png",
             MARG_L, CONT_Y + Cm(0.2), fig_w, max_h=Cm(12.5))

    rx = MARG_L + fig_w + Cm(0.5)
    _add_img(slide, FIGURES / "dist_histogram.png",
             rx, CONT_Y + Cm(0.2), fig_w, max_h=Cm(12.5))

    _add_text(slide,
              "Izq.: boxplots mensuales (todas las estaciones, todos los años)  ·  Der.: distribución global (escala log)",
              MARG_L, SH - Cm(1.3), SW - MARG_L * 2, Cm(0.8),
              size=9, color=GRIS_MED, align=PP_ALIGN.CENTER)


def slide_metodologia_anomalias(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Detección de Anomalías — Marco Metodológico",
               "4 capas independientes · confirmación por acuerdo ≥2 capas")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    capas = [
        ("Capa 1", "Extremos\nUnivariados",
         "Z-score + Hampel ajustado\n(medcouple, |z|>3.5)", "1,634 flags", AZUL_OSC),
        ("Capa 2", "Outliers\nTemporales",
         "Hampel MAD robusta\npor serie de estación", "37,026 flags", AZUL_MED),
        ("Capa 3", "Anomalías\nEspaciales",
         "Kriging CV + LOF\n+ LISA (Moran local)", "14,906 flags", VERDE),
        ("Capa 4", "Anomalías\nMultivariadas",
         "IsoForest + Autoencoder\n+ MCD robusta", "33,780 flags", NARANJA),
    ]

    box_w = (SW - MARG_L * 2 - Cm(0.6)) / 4
    box_h = Cm(8.5)
    bx    = MARG_L
    by    = CONT_Y + Cm(0.5)

    for (lbl, name, desc, flags, col) in capas:
        # Header colored
        _add_rect(slide, bx, by, box_w - Cm(0.15), Cm(1.2), fill=col)
        _add_text(slide, lbl, bx, by + Cm(0.1), box_w - Cm(0.15), Cm(1.0),
                  size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        # Body
        _add_rect(slide, bx, by + Cm(1.2), box_w - Cm(0.15),
                  box_h - Cm(1.2), fill=GRIS_CLAR)
        _add_text(slide, name,
                  bx + Cm(0.15), by + Cm(1.4), box_w - Cm(0.4), Cm(1.5),
                  size=14, bold=True, color=col, align=PP_ALIGN.CENTER, wrap=True)
        _add_text(slide, desc,
                  bx + Cm(0.15), by + Cm(3.1), box_w - Cm(0.4), Cm(2.5),
                  size=10, color=GRIS_OSC, align=PP_ALIGN.CENTER, wrap=True)
        # Flags count
        _add_rect(slide, bx + Cm(0.15), by + box_h - Cm(1.8),
                  box_w - Cm(0.45), Cm(1.5), fill=col)
        _add_text(slide, flags,
                  bx + Cm(0.15), by + box_h - Cm(1.75),
                  box_w - Cm(0.45), Cm(1.4),
                  size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
        bx += box_w

    # Flecha y resultado
    arrow_y = CONT_Y + Cm(0.5) + box_h + Cm(0.4)
    _add_rect(slide, MARG_L, arrow_y, SW - MARG_L * 2, Cm(1.6), fill=AZUL_OSC)
    _add_text(slide, "Catálogo consolidado (≥2 capas):  14,290 anomalías  ·  "
              "Clasificación: artefacto / inconsistencia_espacial / evento_extremo / indeterminado",
              MARG_L + Cm(0.3), arrow_y + Cm(0.2), SW - MARG_L * 2 - Cm(0.6), Cm(1.2),
              size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER, wrap=True)


def slide_anomalias_resultado(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Resultados — Catálogo de Anomalías",
               "14,290 celdas confirmadas · Fleiss κ = 0.065 (complementariedad entre capas)")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    fig_w = SW - MARG_L * 2
    _add_img(slide, FIGURES / "anomaly_consolidation_summary.png",
             MARG_L, CONT_Y + Cm(0.1), fig_w, max_h=Cm(13.0))


def slide_coda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Análisis Composicional (CoDA)",
               "El perfil estacional como composición en el símplex S¹²")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    # Texto izquierda
    tx_w = Cm(12)
    txb  = slide.shapes.add_textbox(MARG_L, CONT_Y + Cm(0.3), tx_w, Cm(13))
    tf   = txb.text_frame
    tf.word_wrap = True

    _add_para(tf, "¿Por qué CoDA?", size=14, bold=True, color=AZUL_OSC)
    _add_para(tf, "La precipitación mensual relativa (proporción del total anual) "
              "es una composición: Σwⱼ = 1, wⱼ ≥ 0. "
              "El análisis convencional en R¹² viola la restricción del símplex "
              "y produce distancias y correlaciones espurias.",
              size=11, color=GRIS_OSC, space_before=4)

    _add_para(tf, "Subconjunto analítico", size=13, bold=True,
              color=AZUL_OSC, space_before=12)
    bullets = [
        "≥ 10 meses con datos por año",
        "≥ 3 años válidos por estación",
        "Coordenadas dentro del polígono de México",
        "634 artefactos recodificados como NaN",
    ]
    for b in bullets:
        _add_para(tf, b, size=11, color=GRIS_OSC, bullet="·")

    _add_para(tf, "Resultado: 1,302 de 1,959 estaciones (66.5%)",
              size=12, bold=True, color=VERDE, space_before=6)

    _add_para(tf, "Tratamiento de ceros", size=13, bold=True,
              color=AZUL_OSC, space_before=12)
    _add_para(tf, "Reemplazo multiplicativo (Martín-Fernández, 2003):\n"
              "δ = 0.65 × mín(w > 0)\nAjuste de no-ceros para mantener Σ = 1",
              size=11, color=GRIS_OSC)

    # Figura derecha
    fx = MARG_L + tx_w + Cm(0.4)
    fw = SW - fx - MARG_R
    _add_img(slide, FIGURES / "coda_subset_selection.png",
             fx, CONT_Y + Cm(0.2), fw, max_h=Cm(13.5))


def slide_ilr(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Transformaciones Log-Ratio",
               "CLR · ILR-SBP climatológico · Isometría verificada")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    fig_w = SW - MARG_L * 2
    _add_img(slide, FIGURES / "coda_logratio_transforms.png",
             MARG_L, CONT_Y + Cm(0.1), fig_w, max_h=Cm(10.5))

    # 3 kpis abajo
    kpis = [
        ("||Ψ Ψᵀ − I||_F = 5.2×10⁻¹⁶", "Matriz de contraste\nortonormal ✓"),
        ("Error isometría ≤ 1.8×10⁻¹⁵", "d_Aitchison ≡ d_Euclídea(ILR) ✓"),
        ("28.1% varianza en ILR-1", "Contraste seco vs húmedo\n(Jun–Oct) domina"),
    ]
    kpi_w  = (SW - MARG_L * 2 - Cm(0.6)) / 3
    kx     = MARG_L
    ky     = SH - Cm(3.2)
    for val, lbl in kpis:
        _add_rect(slide, kx, ky, kpi_w - Cm(0.2), Cm(2.6), fill=AZUL_CLAR)
        _add_text(slide, val, kx + Cm(0.1), ky + Cm(0.2),
                  kpi_w - Cm(0.4), Cm(1.2), size=11, bold=True, color=AZUL_OSC,
                  align=PP_ALIGN.CENTER, wrap=True)
        _add_text(slide, lbl, kx + Cm(0.1), ky + Cm(1.3),
                  kpi_w - Cm(0.4), Cm(1.0), size=9, color=GRIS_OSC,
                  align=PP_ALIGN.CENTER, wrap=True)
        kx += kpi_w


def slide_clustering_k(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Clustering — Diagnóstico de K óptimo",
               "K-Means · Gap statistic · GMM · Bootstrap Jaccard  |  K=2..15")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    fig_w = SW - MARG_L * 2
    _add_img(slide, FIGURES / "clustering_diagnostics.png",
             MARG_L, CONT_Y + Cm(0.1), fig_w, max_h=Cm(10.8))

    # Tabla de criterios
    criterios = [
        ("Silhouette", "K = 2", "Estructura binaria: seco / húmedo"),
        ("Gap (1SE)",  "K = 14", "Gradiente continuo; 1SE activa al final"),
        ("GMM BIC",    "K = 6",  "Mínimo penalizado; modelo parsimonioso"),
        ("Jaccard K=14", "0.63", "Estabilidad moderada (0.60–0.75)"),
    ]
    col_ws = [Cm(4.5), Cm(2.8), Cm(15.5)]
    table_x = MARG_L
    table_y = SH - Cm(3.5)
    table_h = Cm(3.0)

    from pptx.util import Pt as _Pt
    tbl = slide.shapes.add_table(
        len(criterios) + 1, 3, table_x, table_y,
        SW - MARG_L * 2, table_h
    ).table

    # Header
    for ci, (txt, cw) in enumerate(zip(["Criterio", "K óptimo", "Interpretación"], col_ws)):
        tbl.columns[ci].width = cw
        cell = tbl.cell(0, ci)
        cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_OSC
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = BLANCO
                r.font.bold = True
                r.font.size = _Pt(10)
            p.alignment = PP_ALIGN.CENTER

    for ri, (crit, kv, interp) in enumerate(criterios, 1):
        bg = GRIS_CLAR if ri % 2 == 0 else BLANCO
        for ci, txt in enumerate([crit, kv, interp]):
            cell = tbl.cell(ri, ci)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = _Pt(10)
                    r.font.color.rgb = AZUL_OSC if ci == 1 else GRIS_OSC
                    r.font.bold = (ci == 1)
                p.alignment = PP_ALIGN.CENTER if ci < 2 else PP_ALIGN.LEFT


def slide_regimenes_mapa(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Regímenes Pluviométricos — Distribución Geográfica",
               "K=14 clusters sobre 1,302 estaciones  ·  3 métodos comparados")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    fig_w = SW - MARG_L * 2
    _add_img(slide, FIGURES / "regime_maps.png",
             MARG_L, CONT_Y + Cm(0.2), fig_w, max_h=Cm(13.5))


def slide_perfiles(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Perfiles Composicionales por Cluster",
               "Proporción mensual del total anual  ·  Media ± Q25–Q75")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    fig_w = SW - MARG_L * 2
    _add_img(slide, FIGURES / "cluster_profiles_kmeans.png",
             MARG_L, CONT_Y + Cm(0.1), fig_w, max_h=Cm(14.0))


def slide_concordancia(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Concordancia Inter-Método y Etiquetas Climatológicas",
               "ARI / NMI entre K-Means, Ward y GMM  ·  Tipología de regímenes")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    # Heatmap concordancia izquierda
    conc_w = Cm(13)
    _add_img(slide, FIGURES / "method_concordance.png",
             MARG_L, CONT_Y + Cm(0.3), conc_w, max_h=Cm(7))

    # Tabla de regímenes derecha
    rx = MARG_L + conc_w + Cm(0.4)
    rw = SW - rx - MARG_R

    _add_text(slide, "Tipología de regímenes (K-Means)",
              rx, CONT_Y + Cm(0.3), rw, Cm(0.7),
              size=12, bold=True, color=AZUL_OSC)

    from pptx.util import Pt as _Pt
    regimenes = [
        ("lluvias_verano",    "9 clusters", "Pico Jun–Sep, >80% en Jun–Oct"),
        ("bimodal",           "4 clusters", "Picos en Mayo y Ago/Sep"),
        ("lluvias_invierno",  "1 cluster",  "C6: noroeste, 43% en Dic–Feb"),
    ]
    tbl = slide.shapes.add_table(
        len(regimenes) + 1, 3, rx, CONT_Y + Cm(1.2), rw, Cm(3.5)
    ).table
    for ci, (h, cw) in enumerate(
        zip(["Régimen", "N clusters", "Característica"], [rw*0.38, rw*0.22, rw*0.40])
    ):
        tbl.columns[ci].width = int(cw)
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_OSC
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = BLANCO
                r.font.bold = True
                r.font.size = _Pt(9)
            p.alignment = PP_ALIGN.CENTER

    for ri, (reg, nc, carac) in enumerate(regimenes, 1):
        bg = GRIS_CLAR if ri % 2 == 0 else BLANCO
        for ci, txt in enumerate([reg, nc, carac]):
            cell = tbl.cell(ri, ci)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = _Pt(9)
                    r.font.color.rgb = GRIS_OSC
                p.alignment = PP_ALIGN.CENTER if ci == 1 else PP_ALIGN.LEFT

    # Kpis ARI/NMI abajo
    kpis = [
        ("ARI = 0.659", "K-Means ↔ Ward\n(alta concordancia)"),
        ("ARI = 0.493", "K-Means ↔ GMM\n(estructura complementaria)"),
        ("NMI = 0.715", "K-Means ↔ Ward\n(información compartida)"),
    ]
    kpi_w = (SW - MARG_L * 2 - Cm(0.6)) / 3
    kx    = MARG_L
    ky    = SH - Cm(3.8)
    for val, lbl in kpis:
        _add_rect(slide, kx, ky, kpi_w - Cm(0.2), Cm(3.2), fill=AZUL_CLAR)
        _add_text(slide, val, kx + Cm(0.1), ky + Cm(0.3),
                  kpi_w - Cm(0.4), Cm(1.1), size=14, bold=True, color=AZUL_OSC,
                  align=PP_ALIGN.CENTER)
        _add_text(slide, lbl, kx + Cm(0.1), ky + Cm(1.4),
                  kpi_w - Cm(0.4), Cm(1.5), size=9, color=GRIS_OSC,
                  align=PP_ALIGN.CENTER, wrap=True)
        kx += kpi_w


def slide_conclusiones(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "Conclusiones y Trabajo Futuro", "")
    _add_rect(slide, 0, SH - Cm(0.4), SW, Cm(0.4), fill=AZUL_OSC)

    col_w = (SW - MARG_L * 2 - Cm(0.5)) / 2
    rx    = MARG_L + col_w + Cm(0.5)

    # Columna izquierda — Conclusiones
    _add_rect(slide, MARG_L, CONT_Y + Cm(0.2), col_w, Cm(0.7), fill=AZUL_OSC)
    _add_text(slide, "Hallazgos principales",
              MARG_L + Cm(0.2), CONT_Y + Cm(0.2), col_w - Cm(0.4), Cm(0.7),
              size=12, bold=True, color=BLANCO)

    conclusiones = [
        ("Calidad del dato", "46.2% faltante no aleatorio. 634 artefactos identificados "
         "y recodificados antes del análisis composicional."),
        ("Anomalías detectadas", "14,290 celdas confirmadas (4.5% del dataset). "
         "Catálogo clasificado en 4 tipos de anomalía con acción asociada."),
        ("Gradiente climatológico", "El espacio ILR revela un gradiente continuo "
         "dominado por el contraste seco/húmedo (ILR-1, 28% de la varianza)."),
        ("Regímenes identificados", "14 regímenes coherentes geográficamente. "
         "Monzón de verano domina (9/14 clusters). Un cluster de lluvia invernal "
         "en el noroeste. 4 patrones bimodales en el Golfo."),
        ("Robustez metodológica", "K-Means y Ward convergen (ARI=0.66). "
         "Jaccard K=14: 0.63 (moderado), coherente con gradiente continuo."),
    ]

    cy = CONT_Y + Cm(1.1)
    for title, text in conclusiones:
        _add_rect(slide, MARG_L + Cm(0.1), cy, Cm(0.18), Cm(0.5), fill=AZUL_MED)
        _add_text(slide, title,
                  MARG_L + Cm(0.5), cy, col_w - Cm(0.6), Cm(0.55),
                  size=10, bold=True, color=AZUL_OSC)
        cy += Cm(0.55)
        _add_text(slide, text,
                  MARG_L + Cm(0.5), cy, col_w - Cm(0.6), Cm(1.3),
                  size=9, color=GRIS_OSC, wrap=True)
        cy += Cm(1.5)

    # Columna derecha — Trabajo futuro
    _add_rect(slide, rx, CONT_Y + Cm(0.2), col_w, Cm(0.7), fill=NARANJA)
    _add_text(slide, "Trabajo futuro",
              rx + Cm(0.2), CONT_Y + Cm(0.2), col_w - Cm(0.4), Cm(0.7),
              size=12, bold=True, color=BLANCO)

    futuro = [
        ("T3.5.5", "Contraste con Köppen-Geiger (Beck et al., 2023) "
         "y regiones hidrológicas CONAGUA via sjoin espacial + V de Cramér."),
        ("T3.5.6", "Clustering año-por-año (Opción B) para evaluar "
         "estabilidad temporal de los regímenes."),
        ("Tendencias", "¿Se desplaza el inicio de la estación lluviosa "
         "entre 2013 y 2026? Análisis de cambio en la composición año a año."),
        ("Covariables", "Incorporar ENSO (ONI), PDO y AMO como covariables "
         "explicativas de la asignación de cluster."),
    ]

    fy = CONT_Y + Cm(1.1)
    for tag, text in futuro:
        _add_rect(slide, rx + Cm(0.1), fy, Cm(0.18), Cm(0.5), fill=NARANJA)
        _add_text(slide, tag,
                  rx + Cm(0.5), fy, col_w - Cm(0.6), Cm(0.55),
                  size=10, bold=True, color=NARANJA)
        fy += Cm(0.55)
        _add_text(slide, text,
                  rx + Cm(0.5), fy, col_w - Cm(0.6), Cm(1.3),
                  size=9, color=GRIS_OSC, wrap=True)
        fy += Cm(1.5)


def slide_cierre(prs: Presentation, today: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SW, SH, fill=AZUL_OSC)
    _add_rect(slide, 0, SH * 0.7, SW, SH * 0.3, fill=AZUL_MED)
    _add_rect(slide, 0, SH * 0.7, SW, Cm(0.15), fill=BLANCO)

    _add_text(slide, "Gracias",
              Cm(2), Cm(2.5), SW - Cm(4), Cm(3.5),
              size=52, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    _add_text(slide, "Análisis de Precipitación Atípica en México\n"
              "Detección de Anomalías · CoDA · Regímenes Pluviométricos",
              Cm(2), Cm(5.5), SW - Cm(4), Cm(2.5),
              size=15, color=AZUL_CLAR, align=PP_ALIGN.CENTER, wrap=True)

    _add_text(slide, f"UNAM · DCIC · {today}",
              Cm(2), SH * 0.7 + Cm(0.5), SW - Cm(4), Cm(1),
              size=13, color=BLANCO, align=PP_ALIGN.CENTER)
    _add_text(slide, "Código fuente: src/  ·  Reporte completo: outputs/reports/",
              Cm(2), SH * 0.7 + Cm(1.4), SW - Cm(4), Cm(0.8),
              size=10, color=AZUL_CLAR, align=PP_ALIGN.CENTER)


# ── Punto de entrada ──────────────────────────────────────────────────────────

def run_slides(verbose: bool = True) -> Path:
    """Genera el PPTX completo."""
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    today = date.today().strftime("%d de %B de %Y")

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slides_fn = [
        ("Portada",           lambda p: slide_portada(p, today)),
        ("Agenda",            slide_agenda),
        ("Dataset",           slide_dataset),
        ("Datos faltantes",   slide_faltantes),
        ("Distribución",      slide_distribucion),
        ("Metodología anomalías", slide_metodologia_anomalias),
        ("Resultados anomalías",  slide_anomalias_resultado),
        ("CoDA",              slide_coda),
        ("ILR",               slide_ilr),
        ("Clustering K",      slide_clustering_k),
        ("Mapa regímenes",    slide_regimenes_mapa),
        ("Perfiles",          slide_perfiles),
        ("Concordancia",      slide_concordancia),
        ("Conclusiones",      slide_conclusiones),
        ("Cierre",            lambda p: slide_cierre(p, today)),
    ]

    for name, fn in slides_fn:
        if verbose:
            print(f"  Diapositiva: {name}...")
        fn(prs)

    prs.save(str(OUTPUT_PPTX))
    size_kb = OUTPUT_PPTX.stat().st_size / 1024
    if verbose:
        print(f"\n[OK] PPTX generado: {OUTPUT_PPTX}")
        print(f"     {len(slides_fn)} diapositivas  ·  {size_kb:.0f} KB")
    return OUTPUT_PPTX


if __name__ == "__main__":
    run_slides(verbose=True)
